
import streamlit as st
from datetime import datetime
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import openai
import os

openai.api_key = st.secrets["OPENAI_API_KEY"]

st.set_page_config(page_title="MFL Lesson Generator", layout="wide")
st.title("Modern Foreign Language Lesson Generator")

# User inputs
language = st.selectbox("Select Language", ["Spanish", "French", "German"])
year = st.selectbox("Select Year Group", ["7", "8", "9", "10", "11"])
exam_board = st.selectbox("Select Exam Board", ["AQA", "Edexcel", "OCR"])
topic = st.text_input("Topic (e.g., Mi casa)")
focus = st.multiselect("Thematic Focus", ["Grammar", "Vocabulary", "Listening", "Speaking", "Reading", "Writing"])
output_type = "Full lesson pack"

def generate_prompt(language, year, exam_board, topic, focus):
    return f"""
Create a [{output_type}] for [{language}], aimed at Year {year} students following the [{exam_board}] curriculum.
Topic: ["{topic}"]
Thematic Focus: [{', '.join(focus)}]

Ensure the following structure:

1. Title Slide:
   - Auto-filled current date
   - Topic-themed title in the target language
   - Starter task (~8 min): include vocabulary, translation, and tiered challenge levels

2. For Years 7–9:
   - 2–4 vocabulary slides (1–2 words each) with click-to-reveal-style animations
   - Challenge tasks: fill in the blanks, conjugation, or vocab image matching
   - Include engaging visuals

3. For Years 10–11:
   - Reading, listening/dictation, and translation activities
   - Tasks should increase in complexity with longer texts and image analysis

4. Worksheet:
   - Scaffolded vocabulary and sentence-building
   - Extra challenges and visual prompts
   - Include at least one game (e.g., 'one pen, one dice')
"""

def generate_docx(content, filename):
    doc = Document()
    doc.add_heading(f"{language} Worksheet: {topic}", 0)
    doc.add_paragraph(content)
    path = f"/mnt/data/{filename}.docx"
    doc.save(path)
    return path

def generate_pptx(content, filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{language} - {topic}"
    subtitle.text = f"Generated lesson content:\n\n{content[:200]}..."
    path = f"/mnt/data/{filename}.pptx"
    prs.save(path)
    return path

if st.button("Generate Lesson Pack"):
    with st.spinner("Generating lesson content..."):
        prompt = generate_prompt(language, year, exam_board, topic, focus)

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )

        lesson_text = response['choices'][0]['message']['content']
        st.markdown("### Generated Lesson Plan:")
        st.write(lesson_text)

        docx_path = generate_docx(lesson_text, f"{language}_Y{year}_{topic}_Worksheet")
        pptx_path = generate_pptx(lesson_text, f"{language}_Y{year}_{topic}_Lesson")

        st.success("Download your resources below:")
        st.download_button("Download Worksheet (.docx)", data=open(docx_path, "rb"), file_name=os.path.basename(docx_path))
        st.download_button("Download Presentation (.pptx)", data=open(pptx_path, "rb"), file_name=os.path.basename(pptx_path))
