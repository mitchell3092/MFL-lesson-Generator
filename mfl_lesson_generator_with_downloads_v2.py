
import streamlit as st
from io import BytesIO
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def generate_docx(topic, language, year, exam_board, focus):
    doc = Document()
    doc.add_heading(f"{language} Worksheet: {topic}", 0)
    doc.add_paragraph(f"Year Group: {year}")
    doc.add_paragraph(f"Exam Board: {exam_board}")
    doc.add_paragraph("Focus: " + ", ".join(focus))
    doc.add_paragraph("\nTask 1: Match the Spanish to the English.")
    doc.add_paragraph("1. jugar al fútbol  -  a. to play football")
    doc.add_paragraph("2. nadar  -  b. to swim")
    doc.add_paragraph("3. correr  -  c. to run")
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    return doc_io

def generate_pptx(topic, language, year, exam_board, focus):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = f"{language} Lesson: {topic}"

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf_c = content.text_frame
    tf_c.text = f"Year: {year}\nExam Board: {exam_board}\nFocus: {', '.join(focus)}\n\nVocabulary:"

    for item in ["jugar al fútbol - to play football", "nadar - to swim", "correr - to run"]:
        p = tf_c.add_paragraph()
        p.text = item
        p.font.size = Pt(20)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

st.set_page_config(page_title="MFL Lesson Generator", page_icon="📚")

st.title("MFL Lesson Generator")
st.markdown("Generate lesson resources for MFL (Spanish, French, German).")

language = st.selectbox("Language", ["Spanish", "French", "German"])
year = st.selectbox("Year Group", ["Year 7", "Year 8", "Year 9", "Year 10", "Year 11"])
exam_board = st.selectbox("Exam Board", ["AQA", "Edexcel", "OCR"])
topic = st.text_input("Topic (e.g., Mi familia)", placeholder="Enter your lesson topic")
focus = st.multiselect("Thematic Focus", ["Vocabulary", "Grammar", "Listening", "Speaking", "Reading", "Writing"])

if st.button("Generate Resources"):
    if topic.strip() == "":
        st.warning("Please enter a topic.")
    else:
        st.success("Prompt and files ready!")

        generated_prompt = (
            f"Create a [Full lesson pack] for [{language}], aimed at {year} students "
            f"following the [{exam_board}] curriculum.\n"
            f"Topic: [\"{topic}\"]\n"
            f"Thematic Focus: [{' / '.join(focus)}]\n"
            f"Include scaffolded tasks, challenge options, visuals, and interactive elements."
        )

        st.text_area("Generated Prompt", value=generated_prompt, height=250)

        docx_data = generate_docx(topic, language, year, exam_board, focus)
        st.download_button("Download Worksheet (.docx)", docx_data, file_name="worksheet.docx")

        pptx_data = generate_pptx(topic, language, year, exam_board, focus)
        st.download_button("Download Presentation (.pptx)", pptx_data, file_name="lesson.pptx")
        streamlit
python-docx
python-pptx
